from flask import Flask, request, jsonify
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import TextFormatter
from flask_cors import CORS 
import re
import json
from langchain.llms import GPT4All  
from langchain_groq import ChatGroq
import os
from dotenv import load_dotenv
load_dotenv()
from pymongo import MongoClient
import google.generativeai as genai
import io
from google.generativeai import GenerativeModel
import jwt
from functools import wraps
from werkzeug.utils import secure_filename
import logging
from bson.objectid import ObjectId
from langchain_core.prompts import ChatPromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import PyPDFLoader
from langchain.indexes import VectorstoreIndexCreator
from langchain.vectorstores import FAISS
from langchain.embeddings import SentenceTransformerEmbeddings
from langchain.document_loaders import PyMuPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import RetrievalQA
from io import BytesIO
from PyPDF2 import PdfReader  
from langchain.schema import Document  
import chromadb
from sentence_transformers import SentenceTransformer
import google.generativeai as genai
from langchain.document_loaders import PyPDFLoader
from pptx import Presentation

app = Flask(__name__)
 
app = Flask(__name__)
SECRET_KEY = "quick" 
mongo_client = MongoClient("mongodb://localhost:27017/") 
db = mongo_client["quicklearnai"]
topics_collection = db["statistics"]

CORS(app, resources={
    r"/*": {
        "origins": ["http://localhost:5173", "http://localhost:3000"],
        "methods": ["GET", "POST", "OPTIONS"],
        "allow_headers": ["Content-Type", "Authorization"],
        "supports_credentials": True
    }
})

formatter = TextFormatter()

def get_and_enhance_transcript(youtube_url):
    try:
        video_id = youtube_url.split('v=')[-1]
        transcript = None
        language = None

        try:
            transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['hi'])
            language = 'hi'
        except:
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['en'])
                language = 'en'
            except:
                return None, None

        formatted_transcript = formatter.format_transcript(transcript)

        prompt = f"""
        Act as a transcript cleaner. Generate a new transcript with the same context and the content only covered in the given transcript. 
        If there is a revision portion, differentiate it with the actual transcript.
        Give the results in sentences line by line, not in a single line. Also check whether the transcript words have any educational content relevance or not; if not then just give output as: 'Fake transcript'.
        Transcript: {formatted_transcript}
        """
        # apikey = os.getenv("GROQ_API_KEY")
        llm = ChatGroq(
            model="llama-3.3-70b-specdec",
            temperature=0,
            groq_api_key=os.getenv("GROQ_API_KEY")
        )

        enhanced_transcript = llm.invoke(prompt)

        return enhanced_transcript, language
    except Exception as e:
        print(f"Error: {str(e)}")
        return None, None

def generate_summary_and_quiz(transcript, num_questions, language, difficulty):

    try:
        
        prompt = f"""
     
        Summarize the following transcript by identifying the key topics covered, and provide a detailed summary of each topic in 6-7 sentences.
        Each topic should be labeled clearly as "Topic X", where X is the topic name. Provide the full summary for each topic in English, even if the transcript is in a different language.
        Strictly ensure that possessives (e.g., John's book) and contractions (e.g., don't) use apostrophes (') instead of quotation marks (" or “ ”).

        If the transcript contains 'Fake Transcript', do not generate any quiz or summary.

        After the summary, give the name of the topic on which the transcript was all about in a maximum of 2 to 3 words.
        After summarizing, create a quiz with {num_questions} multiple-choice questions in English, based on the transcript content.
        Only generate {difficulty} difficulty questions. Format the output in JSON format as follows, just give the JSON as output, nothing before it:

        {{
            "summary": {{
                "topic1": "value1",
                "topic2": "value2",
                "topic3": "value3"
            }},
            "questions": {{
                "{difficulty}": [
                    {{
                        "question": "What is the capital of France?",
                        "options": ["Paris", "London", "Berlin", "Madrid"],
                        "answer": "Paris"
                    }},
                    {{
                        "question": "What is the capital of Germany?",
                        "options": ["Paris", "London", "Berlin", "Madrid"],
                        "answer": "Berlin"
                    }}
                ]
            }}
        }}

        Transcript: {transcript}
        """
        llm = ChatGroq(
            model="llama-3.3-70b-specdec",
            temperature=0,
            groq_api_key=os.getenv("GROQ_API_KEY")
        )
        response = llm.invoke(prompt)
        if hasattr(response, 'content'):
            response_content = response.content
            try:
                print("response_content:",response_content)
                response_dict = json.loads(response_content)
                print("response_dict:",response_dict)
                return response_dict
            except json.JSONDecodeError as e:
                print(f"JSONDecodeError: {e}")
                return None
        else:
            print("Response does not have a 'content' attribute.")
            return None

    except Exception as e:
        print(f"Error generating summary and quiz: {str(e)}")
        return None

@app.route('/quiz', methods=['POST', 'OPTIONS'])
def quiz():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
        
    data = request.json
    youtube_link = data.get('link')
    num_questions = data.get('qno')
    difficulty = data.get('difficulty')

    if youtube_link:
        transcript, language = get_and_enhance_transcript(youtube_link)
        
        if transcript:
            summary_and_quiz = generate_summary_and_quiz(transcript, num_questions, language, difficulty)
            
            if summary_and_quiz:
                print(summary_and_quiz) 
                return jsonify(summary_and_quiz)

            else:
                return jsonify({"error": "Failed to generate quiz"}), 500
        else:
            return jsonify({"error": "Failed to fetch transcript"}), 404
    else:
        return jsonify({"error": "No YouTube URL provided"}), 400
    


# recommendation
def validate_token_middleware():
    def middleware(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            auth_header = request.headers.get("Authorization")
            token = auth_header.split("Bearer ")[-1] if auth_header and "Bearer " in auth_header else None
            
            if not token:
                return jsonify({"message": "Unauthorized: No token provided"}), 401
            
            try:
                decoded = jwt.decode(token, SECRET_KEY, algorithms=["HS256"])
                request.user_id = decoded.get("id")
                request.user_role = decoded.get("role")  # Optional
                return func(*args, **kwargs)
            except jwt.ExpiredSignatureError:
                return jsonify({"message": "Unauthorized: Token has expired"}), 401
            except jwt.InvalidTokenError as e:
                print(f"Token decoding error: {e}")
                return jsonify({"message": "Unauthorized: Invalid token"}), 401
        
        return wrapper
    return middleware


def llama_generate_recommendations(prompt):
    try:
        api_key=os.getenv("GENAI_API_KEY")
        genai.configure(api_key=api_key)
        
        model = GenerativeModel('gemini-2.0-flash-exp')
        
        response = model.generate_content(prompt)
        
        return response.text
    except Exception as e:
        return f"Error connecting to Gemini API: {e}"
    
@app.route('/getonly', methods=['GET'])
@validate_token_middleware()
def get_recommendations():
    user_id = request.user_id  # Extract user ID from the token
    try:
        user_documents = topics_collection.find({"student": ObjectId(user_id)})
        user_list = list(user_documents)

        topics = [doc.get("topic") for doc in user_list if "topic" in doc]

        if not topics:
            return jsonify({"message": "No topics found for the provided user."}), 404

        prompt = f"Act as an intelligent recommendation generator. Based on the topics provided, generate a brief yet informative overview for each topic and recommend relevant content. Additionally, provide five working YouTube video URLs for each topic that offer valuable insights, explanations, or tutorials. Ensure that the recommendations are diverse, covering different perspectives, and that the video links are accessible and relevant. " \
                 f"The topics are: {', '.join(topics)}"
        recommendations = llama_generate_recommendations(prompt)

        return jsonify({
            "message": "Recommendations generated successfully",
            "recommendations": recommendations
        }), 200

    except Exception as e:
        print("Error:", str(e))
        return jsonify({"message": f"An error occurred: {str(e)}"}), 500


import faiss 
from sentence_transformers import SentenceTransformer
from huggingface_hub import login
groq_api_key = os.getenv("GROQ_API_KEY")
groq_model_name = "llama3-8b-8192"
login(token=os.getenv("HUGGINGFACE_TOKEN")) 
groq_chat = ChatGroq(
    groq_api_key=groq_api_key,
    model_name=groq_model_name,
)


groq_sys_prompt = ChatPromptTemplate.from_template(
    "You are very smart at everything, you always give the best, the most accurate and most precise answers. "
    "Answer the following questions: {user_prompt}. Add more information as per your knowledge so that user can get proper knowledge, but make sure information is correct"
)

embedding_model = SentenceTransformer('multi-qa-mpnet-base-cos-v1')  # Pre-trained model for embeddings
dimension = embedding_model.get_sentence_embedding_dimension()
faiss_index = faiss.IndexFlatL2(dimension) 
metadata_store = {}
pdf_storage = {}

def store_in_faiss(filename, text):
    chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
    embeddings = embedding_model.encode(chunks)
    faiss_index.add(embeddings)  
    metadata_store.update({i: filename for i in range(len(metadata_store), len(metadata_store) + len(chunks))})



genai.configure(api_key=os.getenv("GENAI_API_KEY"))
model = SentenceTransformer("multi-qa-mpnet-base-cos-v1")
chroma_client = chromadb.PersistentClient(path="./chroma_db")
collection = chroma_client.get_or_create_collection(name="pdf_documents")

def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    return " ".join(page.extract_text() for page in reader.pages if page.extract_text())

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):  
                text.append(shape.text)
    return " ".join(text)  

@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file part in the request"}), 400
    
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400
    
    file_ext = os.path.splitext(file.filename)[-1].lower()
    file_path = os.path.join("./uploads", file.filename)
    os.makedirs("./uploads", exist_ok=True)
    file.save(file_path)
    
    try:
        if file_ext == ".pdf":
            content = extract_text_from_pdf(file_path)
        elif file_ext == ".pptx":
            content = extract_text_from_pptx(file_path)
        else:
            return jsonify({"error": "Unsupported file format. Only PDF and PPTX are allowed."}), 400
        
        embedding = model.encode(content).tolist()
        collection.add(documents=[content], embeddings=[embedding], ids=[file.filename])
        
        return jsonify({"message": "File uploaded and processed successfully."}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    
    
@app.route("/query", methods=["POST"])
def query_file():
    data = request.get_json()
    query = data.get("query", "")
    query_embedding = model.encode(query).tolist()
    results = collection.query(query_embeddings=[query_embedding], n_results=3)
    retrieved_texts = "\n".join(results["documents"][0])
    response = genai.GenerativeModel("gemini-1.5-flash").generate_content(retrieved_texts + "\nQuestion: " + query)
    return jsonify({"answer": response.text})


@app.route('/', methods=['GET'])
def health():
    return jsonify({"status": "ok"}) 

if __name__ == '__main__':
    app.run(debug=True, port=5000)
